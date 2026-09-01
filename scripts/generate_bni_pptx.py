#!/usr/bin/env python3
"""Gera BNI_Editavel.pptx (WIDE 13.333x7.5) com objetos nativos python-pptx.

Dependencias: python-pptx, Pillow
Uso: python3 scripts/generate_bni_pptx.py
"""
from __future__ import annotations

import sys
from io import BytesIO
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
IMAGES = ROOT / "images"
OUT_PATH = ROOT / "BNI_Editavel.pptx"

# Layout WIDE
SLIDE_W = Emu(12_192_000)  # 13.333"
SLIDE_H = Emu(6_858_000)  # 7.5"
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
LIGHT_BLUE = RGBColor(0x5B, 0x7F, 0xC7)
STEEL = RGBColor(0x6B, 0x7B, 0x9E)
OFF_WHITE = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MID_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
DARK_GRAY = RGBColor(0x4A, 0x55, 0x68)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF1)


def _set_run_font(
    run,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    italic: bool = False,
) -> None:
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set("typeface", font_name)


def set_char_spacing(paragraph, hundredths_pt: int) -> None:
    for run in paragraph.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(hundredths_pt))


def rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return (int(color[0]), int(color[1]), int(color[2]))


def faded_bg_png(src: Path, bg: RGBColor, opacity: float = 0.15) -> BytesIO:
    """Original cobrindo 13.333x7.5 (cover) a ~15% sobre a cor de fundo."""
    w, h = 1920, 1080
    img = Image.open(src).convert("RGBA")
    scale = max(w / img.width, h / img.height)
    nw, nh = int(img.width * scale), int(img.height * scale)
    img = img.resize((nw, nh), Image.Resampling.LANCZOS)
    left, top = (nw - w) // 2, (nh - h) // 2
    img = img.crop((left, top, left + w, top + h))
    base = Image.new("RGBA", (w, h), rgb_tuple(bg) + (255,))
    blended = Image.blend(base, img, opacity)
    buf = BytesIO()
    blended.convert("RGB").save(buf, format="PNG", optimize=True)
    buf.seek(0)
    return buf


def zero_margins(tf) -> None:
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    font_name: str = "Calibri",
    font_size: int = 12,
    color: RGBColor = NAVY,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    spacing: int | None = None,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    txbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    tf.anchor = anchor
    zero_margins(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run_font(run, font_name, font_size, color, bold, italic)
    if spacing is not None:
        set_char_spacing(p, spacing)
    return txbox


def add_multiline(
    slide,
    lines: list[tuple[str, dict]],
    x: float,
    y: float,
    w: float,
    h: float,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    """lines: (text, font kwargs) — kwargs aceita font_name/font_size/color/bold/italic/spacing."""
    txbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.anchor = anchor
    zero_margins(tf)
    for i, (text, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = kw.get("align", align)
        p.space_before = Pt(kw.get("space_before", 0))
        p.space_after = Pt(kw.get("space_after", 0))
        run = p.add_run()
        run.text = text
        _set_run_font(
            run,
            kw.get("font_name", "Calibri"),
            kw.get("font_size", 12),
            kw.get("color", NAVY),
            kw.get("bold", False),
            kw.get("italic", False),
        )
        if "spacing" in kw:
            set_char_spacing(p, kw["spacing"])
    return txbox


def no_line(shape) -> None:
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, fill, rounded: bool = False, adj: float = 0.15):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    no_line(shape)
    if rounded:
        try:
            shape.adjustments[0] = adj
        except (IndexError, ValueError):
            pass
    return shape


def fill_shape_text(
    shape,
    text: str,
    font_name: str,
    font_size: int,
    color: RGBColor,
    bold: bool = True,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.CENTER,
) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    tf.anchor = MSO_ANCHOR.MIDDLE
    zero_margins(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    _set_run_font(run, font_name, font_size, color, bold, italic)


def add_badge(slide, number: str, x: float = 0.4, y: float = 0.3) -> None:
    shape = add_rect(slide, x, y, 0.9, 0.9, NAVY, rounded=True, adj=0.18)
    fill_shape_text(shape, number, "Arial Black", 28, WHITE, bold=True)


def add_circle_icon(
    slide,
    text_char: str,
    x: float,
    y: float,
    size: float = 0.45,
    font_size: int = 13,
    fill=None,
    glyph=None,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or NAVY
    no_line(shape)
    tf = shape.text_frame
    tf.word_wrap = False
    tf.anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text_char
    font = "Segoe UI Emoji" if ord(text_char[0]) > 0xFFFF else "Calibri"
    _set_run_font(run, font, font_size, glyph or WHITE, bold=True)
    return shape


def add_p_mark(slide, x: float, y: float, size: float, fill, glyph_color: RGBColor):
    shape = add_rect(slide, x, y, size, size, fill, rounded=True, adj=0.22)
    fill_shape_text(shape, "P", "Arial Black", int(22 * size / 0.55), glyph_color, bold=True)
    return shape


def add_logo(slide, x: float, y: float, dark_bg: bool) -> None:
    """Logo Plano's: quadrado P + wordmark."""
    box_fill = WHITE if dark_bg else NAVY
    glyph = NAVY if dark_bg else WHITE
    title_c = WHITE if dark_bg else NAVY
    add_p_mark(slide, x, y, 0.55, box_fill, glyph)
    add_text(
        slide,
        "PLANO'S",
        x + 0.68,
        y - 0.06,
        2.35,
        0.38,
        font_name="Arial Black",
        font_size=20,
        color=title_c,
        bold=True,
        word_wrap=False,
    )
    add_text(
        slide,
        "CONTABILIDADE",
        x + 0.68,
        y + 0.28,
        2.40,
        0.28,
        font_name="Arial",
        font_size=10,
        color=MID_GRAY,
        bold=False,
        word_wrap=False,
        spacing=300,
    )


def add_bg_image(slide, img_path: Path, overlay_color: RGBColor) -> None:
    """Insere a imagem original full-slide a ~15% (blend sobre a cor de fundo)."""
    buf = faded_bg_png(img_path, overlay_color, opacity=0.15)
    slide.shapes.add_picture(buf, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))


def add_icon_item(
    slide,
    icon: str,
    title: str,
    desc: str | None,
    x: float,
    y: float,
    text_w: float,
    title_size: int = 12,
    desc_size: int = 10,
    title_color: RGBColor = NAVY,
    desc_color: RGBColor = STEEL,
    desc_italic: bool = False,
    icon_size: float = 0.45,
    icon_font: int = 13,
) -> None:
    add_circle_icon(slide, icon, x, y, icon_size, icon_font)
    tx = x + icon_size + 0.16
    add_text(
        slide,
        title,
        tx,
        y - 0.02,
        text_w,
        0.32 if desc else 0.50,
        font_name="Calibri",
        font_size=title_size,
        color=title_color,
        bold=True,
        word_wrap=True,
        anchor=MSO_ANCHOR.MIDDLE if not desc else MSO_ANCHOR.TOP,
    )
    if desc:
        add_text(
            slide,
            desc,
            tx,
            y + 0.26,
            text_w,
            0.32,
            font_name="Calibri",
            font_size=desc_size,
            color=desc_color,
            italic=desc_italic,
            word_wrap=True,
        )


def blank_slide(prs: Presentation):
    layout = None
    for candidate in prs.slide_layouts:
        if candidate.name.lower() == "blank":
            layout = candidate
            break
    if layout is None:
        layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    return slide


def set_slide_color(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ---------------------------------------------------------------------------
# SLIDES
# ---------------------------------------------------------------------------

def build_slide_1(prs: Presentation) -> None:
    slide = blank_slide(prs)
    set_slide_color(slide, NAVY)
    add_bg_image(slide, IMAGES / "image1.png", NAVY)

    add_logo(slide, 0.6, 0.4, dark_bg=True)

    add_text(slide, "CONTABILIDADE", 0.6, 1.70, 9.5, 0.78,
             "Arial Black", 42, WHITE, bold=True, word_wrap=False)
    add_text(slide, "QUE ORGANIZA,", 0.6, 2.50, 9.5, 0.72,
             "Arial Black", 42, WHITE, bold=True, word_wrap=False)
    add_text(slide, "CUIDA E FAZ SUA", 0.6, 3.25, 10.0, 0.68,
             "Arial Black", 38, LIGHT_BLUE, bold=True, italic=True, word_wrap=False)
    add_text(slide, "EMPRESA CRESCER!", 0.6, 3.95, 10.5, 0.68,
             "Arial Black", 38, LIGHT_BLUE, bold=True, italic=True, word_wrap=False)

    add_multiline(
        slide,
        [
            ("SEGURANÇA PARA EMPREENDER.", {
                "font_name": "Calibri", "font_size": 14, "color": MID_GRAY, "bold": True,
            }),
            ("TRANQUILIDADE PARA DECIDIR.", {
                "font_name": "Calibri", "font_size": 14, "color": MID_GRAY, "bold": True,
                "space_before": 4,
            }),
        ],
        0.6, 5.10, 8.5, 0.70,
    )

    add_text(
        slide,
        "→ APRESENTAÇÃO BNI",
        0.6, 6.50, 7.0, 0.40,
        "Calibri", 12, WHITE, bold=True, word_wrap=False, spacing=280,
    )


def build_slide_2(prs: Presentation) -> None:
    slide = blank_slide(prs)
    set_slide_color(slide, OFF_WHITE)
    add_bg_image(slide, IMAGES / "image2.png", OFF_WHITE)

    add_badge(slide, "02", 0.4, 0.3)
    add_logo(slide, 10.15, 0.28, dark_bg=False)

    add_text(slide, "CUIDADOS NA", 0.4, 1.50, 7.2, 0.58,
             "Arial Black", 33, NAVY, bold=True, word_wrap=False)
    add_text(slide, "ABERTURA DE", 0.4, 2.10, 7.2, 0.58,
             "Arial Black", 33, NAVY, bold=True, word_wrap=False)
    add_text(slide, "EMPRESAS", 0.4, 2.70, 7.2, 0.62,
             "Arial Black", 38, NAVY, bold=True, word_wrap=False)

    items = [
        ("⚖", "Escolha correta da natureza jurídica", "Alinhada ao seu objetivo e ao seu momento."),
        ("💰", "Enquadramento tributário adequado", "Evita pagar impostos desnecessários."),
        ("📄", "Regularização completa desde o início", "CNPJ, Inscrições, alvarás e licenças em dia."),
        ("💡", "Planejamento é o melhor caminho", "Mais segurança, menos riscos e economia."),
    ]
    y0 = 3.70
    gap = 0.85
    for i, (icon, title, desc) in enumerate(items):
        add_icon_item(slide, icon, title, desc, 0.4, y0 + i * gap, text_w=6.4,
                      title_size=12, desc_size=10)


def build_slide_3(prs: Presentation) -> None:
    slide = blank_slide(prs)
    set_slide_color(slide, OFF_WHITE)
    add_bg_image(slide, IMAGES / "image3.png", OFF_WHITE)

    add_badge(slide, "03", 0.4, 0.3)
    add_text(slide, "PRINCIPAIS ERROS", 1.50, 0.25, 8.5, 0.55,
             "Arial Black", 29, NAVY, bold=True, word_wrap=False)
    add_text(slide, "DO MEI", 1.50, 0.78, 8.5, 0.62,
             "Arial Black", 40, NAVY, bold=True, word_wrap=False)

    errors = [
        ("Misturar receitas e despesas pessoais", "Compromete o controle e pode gerar problemas fiscais."),
        ("Não emitir notas fiscais", "Pode gerar multas e até o desenquadramento do MEI."),
        ("Não guardar comprovantes", "Essenciais para justificar movimentações."),
        ("Não acompanhar o faturamento", "Ultrapassar o limite pode trazer prejuízos."),
    ]
    y0 = 2.00
    gap = 1.20
    for i, (title, desc) in enumerate(errors):
        add_icon_item(
            slide, "✕", title, desc, 6.35, y0 + i * gap, text_w=6.3,
            title_size=13, desc_size=11, icon_font=14,
        )

    add_p_mark(slide, 0.6, 6.50, 0.50, NAVY, WHITE)


def build_slide_4(prs: Presentation) -> None:
    slide = blank_slide(prs)
    set_slide_color(slide, OFF_WHITE)
    add_bg_image(slide, IMAGES / "image4.png", OFF_WHITE)

    add_badge(slide, "04", 0.4, 0.3)
    add_text(slide, "SEPARE SEMPRE:", 1.50, 0.28, 8.5, 0.50,
             "Arial Black", 28, NAVY, bold=True, word_wrap=False)
    add_text(slide, "PF E PJ", 1.50, 0.80, 8.5, 0.60,
             "Arial Black", 38, NAVY, bold=True, word_wrap=False)

    items = [
        ("🛡", "Não misture as contas pessoais com as da empresa."),
        ("📈", "Mais controle financeiro e clareza nos resultados."),
        ("✓", "Evita problemas com o Fisco e facilita o crescimento."),
    ]
    y0 = 1.80
    gap = 1.30
    for i, (icon, title) in enumerate(items):
        add_icon_item(
            slide, icon, title, None, 0.35, y0 + i * gap, text_w=5.35,
            title_size=12, title_color=DARK_GRAY, icon_font=12,
        )

    add_rect(slide, 6.5, 1.50, 2.8, 3.5, LIGHT_GRAY, rounded=True, adj=0.08)
    add_circle_icon(slide, "👤", 7.40, 2.05, 0.70, 22)
    add_text(slide, "PF", 6.5, 2.90, 2.8, 0.50,
             "Arial Black", 24, NAVY, bold=True, align=PP_ALIGN.CENTER, word_wrap=False)
    add_text(slide, "PESSOAL", 6.5, 3.40, 2.8, 0.40,
             "Calibri", 14, STEEL, align=PP_ALIGN.CENTER, word_wrap=False, spacing=200)

    add_rect(slide, 9.8, 1.50, 2.8, 3.5, NAVY, rounded=True, adj=0.08)
    add_circle_icon(slide, "📁", 10.70, 2.05, 0.70, 20, fill=WHITE, glyph=NAVY)
    add_text(slide, "PJ", 9.8, 2.90, 2.8, 0.50,
             "Arial Black", 24, WHITE, bold=True, align=PP_ALIGN.CENTER, word_wrap=False)
    add_text(slide, "EMPRESA", 9.8, 3.40, 2.8, 0.40,
             "Calibri", 14, MID_GRAY, align=PP_ALIGN.CENTER, word_wrap=False, spacing=200)

    add_rect(slide, 0.3, 5.80, 5.5, 0.08, NAVY, rounded=False)
    add_text(slide, "EMPRESA ORGANIZADA,", 0.3, 6.00, 8.0, 0.42,
             "Arial Black", 20, NAVY, bold=True, word_wrap=False)
    add_text(slide, "DECISÕES MAIS ASSERTIVAS!", 0.3, 6.45, 9.0, 0.40,
             "Arial Black", 16, LIGHT_BLUE, bold=True, italic=True, word_wrap=False)


def build_slide_5(prs: Presentation) -> None:
    slide = blank_slide(prs)
    set_slide_color(slide, OFF_WHITE)
    add_bg_image(slide, IMAGES / "image5.png", OFF_WHITE)

    add_badge(slide, "05", 0.4, 0.3)
    add_text(slide, "NOSSO ONBOARDING:", 1.50, 0.28, 10.0, 0.48,
             "Arial Black", 26, NAVY, bold=True, word_wrap=False)
    add_text(slide, "CUIDADO DESDE O PRIMEIRO DIA", 1.50, 0.78, 10.5, 0.42,
             "Arial Black", 18, LIGHT_BLUE, bold=True, italic=True, word_wrap=False)

    steps = [
        ("💡", "Entendimento do negócio", "Conhecemos você, seu mercado e seus objetivos."),
        ("📁", "Documentação e cadastros", "Coleta e organização de todos os documentos."),
        ("🧮", "Estruturação contábil e fiscal", "Parametrizações corretas para uma gestão segura."),
        ("🤝", "Orientações iniciais", "Apoio para decisões mais inteligentes desde o início."),
        ("👥", "Acompanhamento próximo", "Conte com a gente em todas as etapas da sua jornada."),
    ]
    y0 = 1.55
    gap = 0.88
    for i, (icon, title, desc) in enumerate(steps):
        add_icon_item(
            slide, icon, title, desc, 0.40, y0 + i * gap, text_w=7.3,
            title_size=13, desc_size=11, desc_italic=True, icon_font=12,
        )

    banner = add_rect(slide, 0.3, 6.20, 12.7, 0.90, NAVY, rounded=True, adj=0.12)
    fill_shape_text(
        banner,
        "✓  ONBOARDING COMPLETO = FUNDAÇÃO FORTE PARA O CRESCIMENTO.",
        "Arial Black",
        16,
        WHITE,
        bold=True,
    )


REQUIRED_STRINGS = {
    1: [
        "PLANO'S",
        "CONTABILIDADE",
        "QUE ORGANIZA,",
        "CUIDA E FAZ SUA",
        "EMPRESA CRESCER!",
        "SEGURANÇA PARA EMPREENDER.",
        "TRANQUILIDADE PARA DECIDIR.",
        "→ APRESENTAÇÃO BNI",
    ],
    2: [
        "02",
        "CUIDADOS NA",
        "ABERTURA DE",
        "EMPRESAS",
        "Escolha correta da natureza jurídica",
        "Alinhada ao seu objetivo e ao seu momento.",
        "Enquadramento tributário adequado",
        "Evita pagar impostos desnecessários.",
        "Regularização completa desde o início",
        "CNPJ, Inscrições, alvarás e licenças em dia.",
        "Planejamento é o melhor caminho",
        "Mais segurança, menos riscos e economia.",
    ],
    3: [
        "03",
        "PRINCIPAIS ERROS",
        "DO MEI",
        "Misturar receitas e despesas pessoais",
        "Compromete o controle e pode gerar problemas fiscais.",
        "Não emitir notas fiscais",
        "Pode gerar multas e até o desenquadramento do MEI.",
        "Não guardar comprovantes",
        "Essenciais para justificar movimentações.",
        "Não acompanhar o faturamento",
        "Ultrapassar o limite pode trazer prejuízos.",
    ],
    4: [
        "04",
        "SEPARE SEMPRE:",
        "PF E PJ",
        "Não misture as contas pessoais com as da empresa.",
        "Mais controle financeiro e clareza nos resultados.",
        "Evita problemas com o Fisco e facilita o crescimento.",
        "PF",
        "PESSOAL",
        "PJ",
        "EMPRESA",
        "EMPRESA ORGANIZADA,",
        "DECISÕES MAIS ASSERTIVAS!",
    ],
    5: [
        "05",
        "NOSSO ONBOARDING:",
        "CUIDADO DESDE O PRIMEIRO DIA",
        "Entendimento do negócio",
        "Conhecemos você, seu mercado e seus objetivos.",
        "Documentação e cadastros",
        "Coleta e organização de todos os documentos.",
        "Estruturação contábil e fiscal",
        "Parametrizações corretas para uma gestão segura.",
        "Orientações iniciais",
        "Apoio para decisões mais inteligentes desde o início.",
        "Acompanhamento próximo",
        "Conte com a gente em todas as etapas da sua jornada.",
        "ONBOARDING COMPLETO = FUNDAÇÃO FORTE PARA O CRESCIMENTO.",
    ],
}


def collect_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


def validate(prs: Presentation) -> None:
    assert len(prs.slides) == 5, f"esperados 5 slides, obtidos {len(prs.slides)}"
    assert prs.slide_width == SLIDE_W
    assert prs.slide_height == SLIDE_H
    errors: list[str] = []
    for idx, required in REQUIRED_STRINGS.items():
        blob = collect_text(prs.slides[idx - 1])
        for s in required:
            if s not in blob:
                errors.append(f"slide {idx}: falta {s!r}")
    if errors:
        raise SystemExit("VALIDACAO FALHOU:\n" + "\n".join(errors))


def main() -> None:
    missing = [IMAGES / f"image{i}.png" for i in range(1, 6) if not (IMAGES / f"image{i}.png").exists()]
    if missing:
        raise SystemExit("Imagens ausentes: " + ", ".join(str(p) for p in missing))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)

    validate(prs)
    prs.save(str(OUT_PATH))
    print(f"Arquivo gerado: {OUT_PATH}")
    print("Validacao: 5 slides, textos obrigatorios presentes, layout WIDE.")


if __name__ == "__main__":
    sys.exit(main() or 0)
